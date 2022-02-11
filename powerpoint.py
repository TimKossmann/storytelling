import datetime
from re import I
from string import hexdigits
from turtle import width
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from data_breaches_bar_chart_bubble_plot_actual_year import Charts_DataBreaches
from data_breaches_attack_vectors_treemap import Chart_AttackVectors
from phishing_graphs import Phishing_Graphs
from passwords_wordcloud import Chart_WordCloud



class Powerpoint():
    def __init__(self):
        self.pp = Presentation("Powerpoint_Template.pptx")
        self.dbr = Charts_DataBreaches()
        self.pg = Phishing_Graphs()
        self.atp = Chart_AttackVectors()
        self.wc = Chart_WordCloud()


        pass

    def get_picture_img_layout(self):
        return self.pp.slide_layouts[0]
    
    def create_images(self):
        date = datetime.date.today()
        actual_year = date.strftime("%Y")

        self.dbr.update_bubblechart_by_year(int(actual_year) -1).write_image("fig_bubblechart_pp.png", width=800, height=400)
        self.dbr.create_lineplot(int(actual_year) -1).write_image("fig_lineplot_pp.png", width=800, height=400,scale=1)
        self.dbr.create_table().write_image('fig_table_pp.png')

        self.atp.fig.write_image("treemap_pp.png", width=600, height=300)
        self.atp.create_treemap_mensch().write_image("treemap_mensch_pp.png", width=600, height=300)
        
        self.wc.create_wordcloud(True).save("word_cloud_pp.png", format="png")     
        
        self.pg.get_link_donut(True, False).write_image("phishing_link_pp.png")
        self.pg.get_input_donut(True, False).write_image("phishing_input_pp.png")
        self.pg.get_attach_donut(True, False).write_image("phishing_attach_pp.png")

        self.pg.get_fail_bar('Branche', None, True).write_image("fail_bar_mark_pp.png", width=900, height=800)
        self.pg.get_fail_bar('Abteilung', None, True).write_image("fail_bar_type_name_pp.png", width=900, height=800)


    def create_pp(self):
        #first_slide = pp.slides.add_slide(layout)
        #pic = first_slide.shapes.add_picture("fail_bar_mark.png", Inches(0), Inches(1.5), height=Inches(5.5))

        bubble_side = self.pp.slides.add_slide(self.get_picture_img_layout())
        title = bubble_side.placeholders[0]
        title.text = "Warum brauch ich das?"

        placeholder = bubble_side.placeholders[1]
        print(placeholder.placeholder_format.type)
        placeholder.insert_picture("fig_lineplot_pp.png")

        sub_text = bubble_side.placeholders[2]
        tf = sub_text.text_frame

        p = tf.add_paragraph()
        p.text = 'Kosten die durch Datenpannen entstehen steigen Jahr für Jahr'
        p.level = 0

        p = tf.add_paragraph()
        p.text = 'Durch die digitalisierung werden Daten immer wertvoller'
        p.level = 0

        p = tf.add_paragraph()
        p.text = 'Es werden immer mehr Daten gespeichert'
        p.level = 0
        p = tf.add_paragraph()
        p = tf.add_paragraph()
        p.text = 'Oft können solche Attacken abgewehrt werden!'
        p.level = 0

        bubble_side = self.pp.slides.add_slide(self.get_picture_img_layout())
        title = bubble_side.placeholders[0]
        title.text = "Firmen mit den größten Schäden durch Datenlecks"

        placeholder = bubble_side.placeholders[1]
        print(placeholder.placeholder_format.type)
        placeholder.insert_picture("fig_bubblechart_pp.png")

        sub_text = bubble_side.placeholders[2]
        tf = sub_text.text_frame

        p = tf.add_paragraph()
        p.text = 'Vor allem Software-Firmen sind betroffen'
        p.level = 0

        p = tf.add_paragraph()
        p.text = 'Die Daten der Kunden sind ihr größtes Gut'
        p.level = 0

        p = tf.add_paragraph()
        p = tf.add_paragraph()
        p.text = 'Nicht nur für Softwareunternehmen relevant\n-> Datenpannen ziehen auch Image Schäden nach sich'
        p.level = 0

        title_slide = self.pp.slides.add_slide(self.pp.slide_layouts[5])
        title = title_slide.placeholders[0]
        title.text = "Hackingmethoden"


        hacker_attacks = self.pp.slides.add_slide(self.pp.slide_layouts[2])
        title = hacker_attacks.placeholders[0]
        title.text = "Welche Schwachstelle führt bei Hackern zum Erfolg?"
        placeholder = hacker_attacks.placeholders[1]
        placeholder.insert_picture("treemap_pp.png")
        sub_text = hacker_attacks.placeholders[2]
        tf = sub_text.text_frame

        p = tf.add_paragraph()
        percentage_human = self.atp.get_percentage_human()
        p.text = f'Mensch ist in {percentage_human}% der Fällen die Schwachstelle'
        p.level = 0

        p = tf.add_paragraph()
        p.text = 'Hackerangriffe die rein auf das System abzielen gibt es seltener als gedacht wird'
        p.level = 0

        p = tf.add_paragraph()
        p = tf.add_paragraph()
        p.text = 'Durch Schulungen kann die Schwachstelle Mensch deutlich verbessert werden'
        p.level = 0

        hacker_human_attacks = self.pp.slides.add_slide(self.pp.slide_layouts[2])
        title = hacker_human_attacks.placeholders[0]
        title.text = "Was sind die erfogreichsten Angriffsarten, mit der Schwachstelle Mensch?"
        placeholder = hacker_human_attacks.placeholders[1]
        placeholder.insert_picture("treemap_mensch_pp.png")
        sub_text = hacker_human_attacks.placeholders[2]
        tf = sub_text.text_frame

        p = tf.add_paragraph()
        p.text = 'Die größten Potentiale für Hacker sind die komprimittierten oder schwachen Anmeldedaten von Mitarbeitern und das Versenden von Phishing-Mails'
        p.level = 0

        p = tf.add_paragraph()
        p.text = 'Meistens ist es Unwissenheit, die zu Fehlern führt'
        p.level = 0

        p = tf.add_paragraph()
        p = tf.add_paragraph()
        p.text = 'Sensibilisierung kann hier schon ein wirksames Mittel sein'
        p.level = 0

        title_slide = self.pp.slides.add_slide(self.pp.slide_layouts[5])
        title = title_slide.placeholders[0]
        title.text = "Phishing"

        phishing_slide = self.pp.slides.add_slide(self.pp.slide_layouts[1])
        title = phishing_slide.placeholders[0]
        title.text = "Was wollen Phishing Mails?"
        
        added_img = 0
        added_txt = 0

        img_names = ["phishing_link_pp.png", "phishing_input_pp.png", "phishing_attach_pp.png"]
        img_txt = [self.pg.get_text_for_dounut("link"), self.pg.get_text_for_dounut("input"), self.pg.get_text_for_dounut("attach")]
        for plch in phishing_slide.placeholders:
            print(plch.placeholder_format.type)
            plc_type = str(plch.placeholder_format.type)
            if "PICTURE" in plc_type:
                #placeholder_loop = phishing_slide.placeholders[index+2]
                plch.insert_picture(img_names[added_img])
                added_img += 1
            if "OBJECT" in plc_type:
                tf = plch.text_frame
                tf.clear()
                p = tf.add_paragraph()
                p.text = img_txt[added_txt].replace("<br>", "")
                added_txt += 1

        phishing_branch = self.pp.slides.add_slide(self.pp.slide_layouts[3])
        title = phishing_branch.placeholders[0]
        title.text = "Welche Branchen fallen auf Phishing rein?"
        placeholder = phishing_branch.placeholders[1]
        placeholder.insert_picture("fail_bar_mark_pp.png")
        sub_text = phishing_branch.placeholders[2]
        tf = sub_text.text_frame

        p = tf.add_paragraph()
        p.text = 'Branchen die digital sind'
        p.level = 0

        p = tf.add_paragraph()
        p.text = 'Branchen die interessante Daten haben'
        p.level = 0

        p = tf.add_paragraph()
        p = tf.add_paragraph()
        p.text = "Branchen die beliebte Ziele sind, sollten auch mehr schulen"

        phishing_branch = self.pp.slides.add_slide(self.pp.slide_layouts[3])
        title = phishing_branch.placeholders[0]
        title.text = "Welche Abteilungen sind anfällig für Phishing?"
        placeholder = phishing_branch.placeholders[1]
        placeholder.insert_picture("fail_bar_type_name_pp.png")
        sub_text = phishing_branch.placeholders[2]
        tf = sub_text.text_frame

        p = tf.add_paragraph()
        p.text = 'Abteilungen die viel Kommunikation betreiben, wie Vetrieb und Marketing sollte beobachtet werden'
        p.level = 0

        p = tf.add_paragraph()
        p.text = 'Gebäudemanagement wird nicht so viel geschult, da sie nicht so oft Phishing Mails erhalten'
        p.level = 0


        p = tf.add_paragraph()
        p = tf.add_paragraph()
        p.text = 'Wichtig alle Abteilungen zu schulen, denn eine Kette ist nur so stark, wie ihr schwächstes Glied'
        p.level = 0
        
        title_slide = self.pp.slides.add_slide(self.pp.slide_layouts[5])
        title = title_slide.placeholders[0]
        title.text = "Passwortsicherheit"

        password_slide = self.pp.slides.add_slide(self.pp.slide_layouts[4])
        title = password_slide.placeholders[0]
        title.text = "Was ist ein schlechtes Passwort?"
        placeholder = password_slide.placeholders[1]
        placeholder.insert_picture("word_cloud_pp.png")
        sub_text = password_slide.placeholders[2]
        tf = sub_text.text_frame

        p = tf.add_paragraph()
        p.text = 'Kurze Passwörter\n-> Die Passwortlänge potentiert die Stärke des Passworts'
        p.level = 0

        p = tf.add_paragraph()
        p.text = 'Passwörter mit wenig Zeichenvariationen\n-> Wenn man Sonderzeichen nutzt erhöht sich die Zahl der Möglichen Zeichen direkt um 32'
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = 'Existierende Wörter\n-> Hacker verwenden Wörterbücher um Passwörter zu knacken'
        p.level = 0 

        """
        placeholder = first_slide.placeholders[1]
        print(placeholder.placeholder_format.type)
        placeholder.insert_picture("fail_bar_mark.png")

        placeholder = first_slide.placeholders[0]
        placeholder.text = 'Fail Bar'

        placeholder = first_slide.placeholders[2]
        placeholder.text = 'Fail Bar \nZweiter Text'
        #print(len(first_slide.placeholders))"""

        self.pp.save("Cyber_Security_LaCTiS.pptx")

if __name__ == "__main__":
    pp = Powerpoint()
    pp.create_images()
    pp.create_pp()
